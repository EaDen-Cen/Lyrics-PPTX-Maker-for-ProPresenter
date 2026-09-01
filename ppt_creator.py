from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from config import SETTINGS
from parser import read_lyrics
import os

# ==========================
# 内部转换参数 (勿动)
# ==========================

PPT_SCALE = 0.5

def px_to_inches(px):
    return px / 144

def calculate_text_offset(
        # 画布宽度对垂直偏移量没影响
        # canvas_width,
        canvas_height,
        chinese_size,
        english_size
):

    # ==========================
    # 基础参数
    # 注：重要数据！不要轻易修改！否则可能会出现计算错误！
    # ==========================

    BASE_HEIGHT = 1080

    # 字号影响
    FONT_FACTOR = 0.58

    # 画布高度影响
    HEIGHT_POWER = 1.35

    # 基础比例
    BASE_RATIO = 0.01074


    # ==========================
    # 计算画布缩放
    # ==========================

    height_scale = (
        canvas_height / BASE_HEIGHT
    ) ** HEIGHT_POWER


    canvas_offset = (
        canvas_height
        * BASE_RATIO
        * height_scale
    )


    # ==========================
    # 字号差异
    # ==========================

    font_offset = (
        chinese_size
        -
        english_size
    ) * FONT_FACTOR


    # ==========================
    # 最终偏移
    # ==========================

    offset = (
        canvas_offset
        +
        font_offset
    )


    return round(offset)

def add_text_slide(
        prs,
        chinese,
        english,
        chinese_size,
        english_size,
        uppercase=True
):
    chinese_font_size = chinese_size * PPT_SCALE
    english_font_size = english_size * PPT_SCALE

    slide = prs.slides.add_slide(
        prs.slide_layouts[6]
    )

    # ==========================
    # 添加备注
    # ==========================

    add_slide_notes(
        slide,
        chinese
    )

    slide.background.fill.background()

    x = Inches(
        px_to_inches(
            (
                SETTINGS["background"]["width"]
                -
                SETTINGS["textbox"]["width"]
            )
            / 2
        )
    )
    y = Inches(
        px_to_inches(
            (
                SETTINGS["background"]["height"]
                -
                SETTINGS["textbox"]["height"]
            )
            / 2
        )
    )
    width = Inches(
        px_to_inches(
            SETTINGS["textbox"]["width"]
        )
    )
    height = Inches(
        px_to_inches(
            SETTINGS["textbox"]["height"] / 2
        )
    )

    # 文本居中

    offset = calculate_text_offset(
        SETTINGS["background"]["height"],
        chinese_size,
        english_size
    )

    y += Inches(px_to_inches(offset))

    # 中文文本框
    chinese_box = slide.shapes.add_textbox(
        x,
        y,
        width,
        height
    )


    # 英文文本框
    english_box = slide.shapes.add_textbox(
        x,
        y + height,
        width,
        height
    )


    from pptx.enum.text import MSO_ANCHOR


    # =====================
    # 中文
    # =====================

    chinese_frame = chinese_box.text_frame
    chinese_frame.clear()

    chinese_frame.vertical_anchor = MSO_ANCHOR.BOTTOM

    chinese_frame.margin_left = 0
    chinese_frame.margin_right = 0
    chinese_frame.margin_top = 0
    chinese_frame.margin_bottom = 0


    p = chinese_frame.paragraphs[0]
    p.text = chinese
    p.alignment = PP_ALIGN.CENTER


    font = p.runs[0].font
    font.name = SETTINGS["font"]["chinese"]
    font.size = Pt(chinese_font_size)
    font.bold = False
    font.color.rgb = RGBColor(255,255,255)



    # =====================
    # 英文
    # =====================

    english_frame = english_box.text_frame
    english_frame.clear()

    english_frame.vertical_anchor = MSO_ANCHOR.TOP

    english_frame.margin_left = 0
    english_frame.margin_right = 0
    english_frame.margin_top = 0
    english_frame.margin_bottom = 0


    p2 = english_frame.paragraphs[0]
    if uppercase:
        p2.text = english.upper()
    else:
        p2.text = english
    p2.alignment = PP_ALIGN.CENTER


    font2 = p2.runs[0].font
    font2.name = SETTINGS["font"]["english"]
    font2.size = Pt(english_font_size)
    font2.bold = False
    font2.color.rgb = RGBColor(255,255,255)

def add_slide_notes(slide, note):
    notes_frame = slide.notes_slide.notes_text_frame

    if notes_frame is not None:
        notes_frame.text = note

def add_lyric_slide(prs, chinese, english):

    add_text_slide(
        prs,
        chinese,
        english,
        SETTINGS["lyric"]["chinese_size"],
        SETTINGS["lyric"]["english_size"],
        True
    )


def add_title_slide(prs, chinese, english):

    add_text_slide(
        prs,
        "《" + chinese + "》",
        english,
        SETTINGS["title"]["chinese_size"],
        SETTINGS["title"]["english_size"],
        False
    )

def create_presentation(
        input_file,
        output_file
):
    lyrics = read_lyrics(input_file)

    prs = Presentation()

    # 设置 16:9
    prs.slide_width = Inches(px_to_inches(SETTINGS["background"]["width"]))
    prs.slide_height = Inches(px_to_inches(SETTINGS["background"]["height"]))

    # 标题

    add_title_slide(
        prs,
        lyrics["title_cn"],
        lyrics["title_en"]
    )

    # Lyrics

    for page in lyrics["slides"]:

        add_lyric_slide(
            prs,
            page["chinese"],
            page["english"]
        )

    output_directory = os.path.dirname(output_file)
    if output_directory:
        os.makedirs(output_directory, exist_ok=True)

    prs.save(output_file)


    print("======================")
    print("歌词PPT生成完成")
    print(output_file)
    print("======================")